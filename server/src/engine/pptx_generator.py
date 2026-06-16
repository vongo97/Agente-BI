import os
import tempfile
import logging
import json
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None
    MSO_SHAPE = None

logger = logging.getLogger(__name__)
from src.utils.exporter import export_plotly_to_image, clean_text

def generate_structured_slide_content(text, api_key=None):
    """
    Analiza el hallazgo mediante Gemini para clasificarlo en un layout de presentación óptimo
    y estructurar su contenido. Retorna un diccionario con el tipo de layout y sus datos.
    """
    fallback_data = {
        "layout_type": "CALLOUT_CARDS",
        "bullets": fallback_summarization(text)
    }

    if not api_key:
        return fallback_data

    prompt = f"""
    Eres un Consultor de BI y Diseñador de Presentaciones de nivel McKinsey/Gartner.
    Tu misión es analizar el siguiente texto de un hallazgo de BI y clasificarlo en uno de los 3 layouts de presentación de PowerPoint para maximizar su impacto visual.
    
    Layouts disponibles y sus condiciones de uso:
    1. "KPI_HIGHLIGHT": Se usa SI Y SOLO SI el texto contiene una métrica numérica muy importante destacada (ej: cifras monetarias grandes, tasas de crecimiento %, múltiplos como 4x, etc.) que deba ser el foco principal de la diapositiva.
       Devuelve JSON:
       {{
         "layout_type": "KPI_HIGHLIGHT",
         "kpi_value": "La cifra o porcentaje corto (máximo 8 caracteres, ej: 6.69M, +15%, 4x)",
         "kpi_label": "Nombre corto de la métrica (máximo 30 caracteres, ej: Ventas Retail, Reacción Analítica)",
         "bullets": [
           "Punto explicativo 1: Detalle del contexto y su impacto...",
           "Punto explicativo 2: Consecuencia o recomendación..."
         ]
       }}

    2. "COMPARISON_2COL": Se usa si el texto compara explícitamente dos entidades, empresas, regiones o períodos de tiempo (ej: Superfinanciera vs Supersociedades, Q1 vs Q2, Región A vs Región B).
       Devuelve JSON:
       {{
         "layout_type": "COMPARISON_2COL",
         "col1_title": "Nombre de la primera entidad/período (ej: Superfinanciera)",
         "col1_bullets": [
           "Punto clave 1: Detalle analítico...",
           "Punto clave 2: Consecuencia..."
         ],
         "col2_title": "Nombre de la segunda entidad/período (ej: Supersociedades)",
         "col2_bullets": [
           "Punto clave 1: Detalle analítico...",
           "Punto clave 2: Consecuencia..."
         ]
       }}

    3. "CALLOUT_CARDS": En cualquier otro caso general de hallazgos del mercado o conclusiones sin métricas masivas ni comparaciones directas.
       Devuelve JSON:
       {{
         "layout_type": "CALLOUT_CARDS",
         "bullets": [
           "Concepto Clave 1: Explicación de soporte del hallazgo...",
           "Concepto Clave 2: Explicación de soporte del hallazgo...",
           "Concepto Clave 3: Explicación de soporte..."
         ]
       }}

    REGLAS IMPORTANTES:
    1. Devuelve ÚNICAMENTE el objeto JSON correspondiente. Sin explicaciones, sin introducciones, sin delimitadores de código markdown extraños.
    2. En "bullets" de cualquier sección, escribe en formato "Concepto clave: Explicación..." para favorecer el resaltado.
    3. Cada viñeta/bullet debe ser corto y conciso (máximo 15 palabras).
    4. Garantiza que las claves del JSON coincidan exactamente con la estructura del layout elegido.

    TEXTO DE BI A ANALIZAR:
    {text}
    """
    try:
        from src.engine.bi_analyst import generate_ai_content
        raw_response = generate_ai_content(prompt, api_key, provider="gemini", temperature=0.1, model_level="ANALYTICS")
        
        if raw_response:
            cleaned = raw_response.strip()
            if cleaned.startswith("```"):
                lines = cleaned.split("\n")
                if lines[0].startswith("```json") or lines[0].startswith("```"):
                    cleaned = "\n".join(lines[1:-1]).strip()
                elif cleaned.endswith("```"):
                    cleaned = cleaned[3:-3].strip()
            
            data = json.loads(cleaned)
            if data.get("layout_type") in ["KPI_HIGHLIGHT", "COMPARISON_2COL", "CALLOUT_CARDS"]:
                return data
    except Exception as e:
        logger.warning("Fallo al generar layout estructurado por IA: %s. Usando fallback.", e)
        
    return fallback_data

def summarize_text_for_slide(text, api_key=None):
    """
    Sintetiza el análisis a una lista simple de viñetas para layouts mixtos (con gráficos).
    """
    if not api_key:
        return fallback_summarization(text)

    prompt = f"""
    Eres un Consultor de BI experto en presentaciones ejecutivas de nivel Directivo (McKinsey/Gartner).
    Tu misión es transformar el siguiente análisis detallado en puntos clave de impacto para una sola diapositiva.

    REGLAS ESTRICTAS:
    1. Devuelve un máximo de 3 a 4 puntos clave (bullets).
    2. Cada punto clave debe ser conciso, potente y tener un máximo de 15 palabras.
    3. Destaca la métrica numérica clave y el impacto comercial.
    4. Cada punto debe iniciar con una viñeta simple "-" (guion) y tener la estructura "Concepto clave: Explicación...".
    5. NO incluyas introducciones, títulos, ni notas. Solo las viñetas.

    ANÁLISIS A SINTETIZAR:
    {text}
    """
    try:
        from src.engine.bi_analyst import generate_ai_content
        summary_raw = generate_ai_content(prompt, api_key, provider="gemini", temperature=0.3, model_level="ANALYTICS")
        if summary_raw and "⚠️" not in summary_raw:
            bullets = [b.strip().lstrip("-*• ").strip() for b in summary_raw.split("\n") if b.strip()]
            bullets = [b for b in bullets if len(b) > 5]
            if len(bullets) >= 2:
                return bullets
    except Exception as e:
        logger.warning("Fallo en síntesis inteligente para diapositiva: %s. Usando fallback.", e)

    return fallback_summarization(text)

def fallback_summarization(text):
    """Extrae las oraciones de mayor impacto de manera determinista."""
    import re
    clean_lines = [line.strip() for line in text.split("\n") if line.strip() and not line.strip().startswith("#") and "|" not in line]
    sentences = []
    for line in clean_lines:
        parts = re.split(r'\.\s+', line)
        for p in parts:
            p_clean = p.strip().rstrip(".")
            if len(p_clean) > 20:
                sentences.append(p_clean)
    return sentences[:3]

def create_presentation(title, summary, items, output_path, template_path="templates/template_vektra_general.pptx", api_key=None):
    """
    Genera un archivo .pptx a partir de hallazgos usando una plantilla maestra.
    """
    if Presentation is None:
        return False, "Librería python-pptx no instalada."

    try:
        # Resolver ruta absoluta del template
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        full_template_path = os.path.join(base_dir, template_path)
        
        if os.path.exists(full_template_path):
            logger.debug("Usando plantilla PPTX: %s", os.path.basename(full_template_path))
            prs = Presentation(full_template_path)
        else:
            logger.warning("Plantilla PPTX no encontrada. Usando presentación vacía.")
            prs = Presentation()
            
        # Eliminar diapositivas sobrantes de la plantilla (conservar solo la 0 de portada)
        while len(prs.slides) > 1:
            rId = prs.slides._sldIdLst[1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[1]

        # Forzar tamaño 16:9 Widescreen para dar protagonismo a los gráficos BI
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
            
        # Portada (Modificar primera diapositiva)
        if len(prs.slides) > 0:
            cover = prs.slides[0]
            # Establecer fondo oscuro Slate 900 en la portada
            cover.background.fill.solid()
            cover.background.fill.fore_color.rgb = RGBColor(15, 23, 42)
            
            # Añadir barra lateral decorativa en Azul Vektra
            from pptx.enum.shapes import MSO_SHAPE
            accent_bar = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=0, top=0, width=Inches(0.4), height=prs.slide_height)
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(37, 99, 235)
            accent_bar.line.fill.background() # Quitar contorno
            
            # Truncar resumen si supera los 250 caracteres
            truncated_summary = summary
            if len(truncated_summary) > 250:
                truncated_summary = truncated_summary[:247] + "..."

            if cover.shapes.title:
                cover.shapes.title.text = title
                # Forzar color de fuente a Blanco recorriendo párrafos
                for paragraph in cover.shapes.title.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.name = "Arial"
                    if len(title) > 30:
                        paragraph.font.size = Pt(36)
                        
            # Buscar placeholders para el resumen
            for shape in cover.placeholders:
                if shape.shape_type == 14 and shape != cover.shapes.title: # Placeholder de texto
                    shape.text = truncated_summary
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Arial"
                        paragraph.font.size = Pt(14)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255) # Forzar blanco
                    break
        else:
            # Crear portada por defecto si la plantilla está totalmente vacía
            slide_layout = prs.slide_layouts[0] 
            slide = prs.slides.add_slide(slide_layout)
            
            # Establecer fondo oscuro Slate 900
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(15, 23, 42)
            
            from pptx.enum.shapes import MSO_SHAPE
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=0, top=0, width=Inches(0.4), height=prs.slide_height)
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(37, 99, 235)
            accent_bar.line.fill.background()
            
            truncated_summary = summary
            if len(truncated_summary) > 250:
                truncated_summary = truncated_summary[:247] + "..."

            slide.shapes.title.text = title
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                
            slide.placeholders[1].text = truncated_summary
            for paragraph in slide.placeholders[1].text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Forzar color de fuente a Blanco en toda la portada (Slide 1) para garantizar legibilidad
        if len(prs.slides) > 0:
            cover = prs.slides[0]
            for shape in cover.shapes:
                if shape.has_text_frame and shape.name != "accent_bar":
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Diapositivas de contenido
        for item in items:
            content_text = item.get("content", "")
            fig_data = item.get("fig")
            
            # Usar layout 5 (Title Only) o el 1 (Title and Content)
            layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            
            # Establecer fondo claro Slate 50 en la diapositiva de contenido
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
            
            # Añadir línea decorativa horizontal en Azul Vektra debajo del encabezado
            from pptx.enum.shapes import MSO_SHAPE
            header_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(0.5), top=Inches(1.6), width=Inches(12.33), height=Inches(0.03))
            header_line.fill.solid()
            header_line.fill.fore_color.rgb = RGBColor(37, 99, 235)
            header_line.line.fill.background() # Quitar borde
            
            # Añadir pie de diapositiva discreto de marca corporativa
            footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4))
            tf_footer = footer_box.text_frame
            tf_footer.word_wrap = True
            p_foot = tf_footer.paragraphs[0]
            p_foot.text = "VEKTRA BI ENGINE | DIAGNÓSTICO ESTRATÉGICO"
            p_foot.font.name = "Arial"
            p_foot.font.size = Pt(8)
            p_foot.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
            
            # Extraer y limpiar título del contenido
            lines = content_text.split("\n")
            slide_title = "Hallazgo Estratégico"
            for line in lines:
                if line.strip().startswith("#"):
                    slide_title = clean_text(line.replace("#", "").strip())
                    break
                    
            # Determinar si la diapositiva tiene gráfico
            has_chart = bool(fig_data)
            
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
                # IMPORTANTE: Forzar el tamaño de la fuente, nombre y color Slate 800
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    paragraph.font.name = "Arial"
                    paragraph.font.size = Pt(24)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(30, 41, 59) # Slate Oscuro
                
            # Sintetizar el contenido usando el Sintetizador Ejecutivo Híbrido
            bullets = summarize_text_for_slide(content_text, api_key)
            
            if not has_chart:
                # Obtener la estructura y layout recomendado por IA (o fallback determinista)
                slide_data = generate_structured_slide_content(content_text, api_key)
                layout_type = slide_data.get("layout_type", "CALLOUT_CARDS")
                
                if layout_type == "KPI_HIGHLIGHT":
                    # Layout A: KPI Gigante a la izquierda, tarjetas de soporte a la derecha
                    kpi_val = slide_data.get("kpi_value", "BI")
                    kpi_lbl = slide_data.get("kpi_label", "Métrica de Diagnóstico")
                    side_bullets = slide_data.get("bullets", [])
                    
                    # 1. Tarjeta de KPI (Izquierda)
                    bg_kpi = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(4.5), Inches(4.5))
                    bg_kpi.fill.solid()
                    bg_kpi.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                    bg_kpi.line.fill.background()
                    
                    accent_kpi = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(0.08), Inches(4.5))
                    accent_kpi.fill.solid()
                    accent_kpi.fill.fore_color.rgb = RGBColor(37, 99, 235) # Azul Vektra
                    accent_kpi.line.fill.background()
                    
                    txBox_kpi = slide.shapes.add_textbox(Inches(0.5 + 0.35), Inches(2.0 + 1.2), Inches(4.5 - 0.5), Inches(4.5 - 1.5))
                    tf_kpi = txBox_kpi.text_frame
                    tf_kpi.word_wrap = True
                    tf_kpi.margin_left = Inches(0)
                    tf_kpi.margin_top = Inches(0)
                    tf_kpi.margin_right = Inches(0)
                    tf_kpi.margin_bottom = Inches(0)
                    
                    p_val = tf_kpi.paragraphs[0]
                    p_val.text = kpi_val
                    p_val.font.name = "Arial"
                    p_val.font.size = Pt(54)
                    p_val.font.bold = True
                    p_val.font.color.rgb = RGBColor(37, 99, 235)
                    p_val.space_after = Pt(12)
                    
                    p_lbl = tf_kpi.add_paragraph()
                    p_lbl.text = kpi_lbl
                    p_lbl.font.name = "Arial"
                    p_lbl.font.size = Pt(14)
                    p_lbl.font.bold = True
                    p_lbl.font.color.rgb = RGBColor(30, 41, 59)
                    
                    # 2. Tarjetas de soporte (Derecha)
                    N = max(1, len(side_bullets))
                    if N == 1:
                        card_h = Inches(2.0)
                        gap = Inches(0)
                    elif N == 2:
                        card_h = Inches(1.8)
                        gap = Inches(0.5)
                    else:
                        card_h = Inches(1.2)
                        gap = Inches(0.3)
                        
                    for i, b_text in enumerate(side_bullets):
                        card_top = Inches(2.0) + i * (card_h + gap)
                        bg_side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), card_top, Inches(7.53), card_h)
                        bg_side.fill.solid()
                        bg_side.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                        bg_side.line.fill.background()
                        
                        accent_side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), card_top, Inches(0.08), card_h)
                        accent_side.fill.solid()
                        accent_side.fill.fore_color.rgb = RGBColor(37, 99, 235)
                        accent_side.line.fill.background()
                        
                        txBox_side = slide.shapes.add_textbox(Inches(5.3 + 0.25), card_top + Inches(0.05), Inches(7.53 - 0.4), card_h - Inches(0.1))
                        tf_side = txBox_side.text_frame
                        tf_side.word_wrap = True
                        tf_side.margin_left = Inches(0)
                        tf_side.margin_top = Inches(0)
                        tf_side.margin_right = Inches(0)
                        tf_side.margin_bottom = Inches(0)
                        
                        p_s = tf_side.paragraphs[0]
                        f_size = Pt(13) if N <= 3 else Pt(11.5)
                        p_s.font.name = "Arial"
                        p_s.font.size = f_size
                        
                        b_text_clean = clean_text(b_text).strip()
                        parts = b_text_clean.split(":", 1)
                        if len(parts) == 2 and len(parts[0]) < 40:
                            run_t = p_s.add_run()
                            run_t.text = parts[0].strip() + ": "
                            run_t.font.bold = True
                            run_t.font.size = f_size
                            run_t.font.color.rgb = RGBColor(30, 41, 59)
                            
                            run_d = p_s.add_run()
                            run_d.text = parts[1].strip()
                            run_d.font.bold = False
                            run_d.font.size = f_size
                            run_d.font.color.rgb = RGBColor(71, 85, 105)
                        else:
                            run_a = p_s.add_run()
                            run_a.text = b_text_clean
                            run_a.font.bold = False
                            run_a.font.size = f_size
                            run_a.font.color.rgb = RGBColor(71, 85, 105)
                            
                elif layout_type == "COMPARISON_2COL":
                    # Layout B: Dos Columnas Comparativas Lado a Lado
                    col1_t = slide_data.get("col1_title", "Concepto A")
                    col1_bl = slide_data.get("col1_bullets", [])
                    col2_t = slide_data.get("col2_title", "Concepto B")
                    col2_bl = slide_data.get("col2_bullets", [])
                    
                    # Columna 1 (Izquierda)
                    bg_col1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(5.9), Inches(4.5))
                    bg_col1.fill.solid()
                    bg_col1.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                    bg_col1.line.fill.background()
                    
                    accent_col1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(0.08), Inches(4.5))
                    accent_col1.fill.solid()
                    accent_col1.fill.fore_color.rgb = RGBColor(37, 99, 235)
                    accent_col1.line.fill.background()
                    
                    txBox_col1 = slide.shapes.add_textbox(Inches(0.5 + 0.25), Inches(2.0 + 0.15), Inches(5.9 - 0.4), Inches(4.5 - 0.3))
                    tf_col1 = txBox_col1.text_frame
                    tf_col1.word_wrap = True
                    tf_col1.margin_left = Inches(0)
                    tf_col1.margin_top = Inches(0)
                    tf_col1.margin_right = Inches(0)
                    tf_col1.margin_bottom = Inches(0)
                    
                    p_t1 = tf_col1.paragraphs[0]
                    p_t1.text = col1_t
                    p_t1.font.name = "Arial"
                    p_t1.font.size = Pt(16)
                    p_t1.font.bold = True
                    p_t1.font.color.rgb = RGBColor(30, 41, 59)
                    p_t1.space_after = Pt(12)
                    
                    for b_text in col1_bl:
                        p_b = tf_col1.add_paragraph()
                        p_b.font.name = "Arial"
                        p_b.font.size = Pt(11)
                        p_b.space_after = Pt(6)
                        
                        run_bl = p_b.add_run()
                        run_bl.text = "• "
                        run_bl.font.color.rgb = RGBColor(75, 85, 99)
                        
                        b_text_clean = clean_text(b_text).strip()
                        parts = b_text_clean.split(":", 1)
                        if len(parts) == 2 and len(parts[0]) < 40:
                            run_t = p_b.add_run()
                            run_t.text = parts[0].strip() + ": "
                            run_t.font.bold = True
                            run_t.font.color.rgb = RGBColor(30, 41, 59)
                            
                            run_d = p_b.add_run()
                            run_d.text = parts[1].strip()
                            run_d.font.bold = False
                            run_d.font.color.rgb = RGBColor(71, 85, 105)
                        else:
                            run_d = p_b.add_run()
                            run_d.text = b_text_clean
                            run_d.font.bold = False
                            run_d.font.color.rgb = RGBColor(71, 85, 105)
                            
                    # Columna 2 (Derecha)
                    bg_col2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(2.0), Inches(5.9), Inches(4.5))
                    bg_col2.fill.solid()
                    bg_col2.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                    bg_col2.line.fill.background()
                    
                    accent_col2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(2.0), Inches(0.08), Inches(4.5))
                    accent_col2.fill.solid()
                    accent_col2.fill.fore_color.rgb = RGBColor(37, 99, 235)
                    accent_col2.line.fill.background()
                    
                    txBox_col2 = slide.shapes.add_textbox(Inches(6.93 + 0.25), Inches(2.0 + 0.15), Inches(5.9 - 0.4), Inches(4.5 - 0.3))
                    tf_col2 = txBox_col2.text_frame
                    tf_col2.word_wrap = True
                    tf_col2.margin_left = Inches(0)
                    tf_col2.margin_top = Inches(0)
                    tf_col2.margin_right = Inches(0)
                    tf_col2.margin_bottom = Inches(0)
                    
                    p_t2 = tf_col2.paragraphs[0]
                    p_t2.text = col2_t
                    p_t2.font.name = "Arial"
                    p_t2.font.size = Pt(16)
                    p_t2.font.bold = True
                    p_t2.font.color.rgb = RGBColor(30, 41, 59)
                    p_t2.space_after = Pt(12)
                    
                    for b_text in col2_bl:
                        p_b = tf_col2.add_paragraph()
                        p_b.font.name = "Arial"
                        p_b.font.size = Pt(11)
                        p_b.space_after = Pt(6)
                        
                        run_bl = p_b.add_run()
                        run_bl.text = "• "
                        run_bl.font.color.rgb = RGBColor(75, 85, 99)
                        
                        b_text_clean = clean_text(b_text).strip()
                        parts = b_text_clean.split(":", 1)
                        if len(parts) == 2 and len(parts[0]) < 40:
                            run_t = p_b.add_run()
                            run_t.text = parts[0].strip() + ": "
                            run_t.font.bold = True
                            run_t.font.color.rgb = RGBColor(30, 41, 59)
                            
                            run_d = p_b.add_run()
                            run_d.text = parts[1].strip()
                            run_d.font.bold = False
                            run_d.font.color.rgb = RGBColor(71, 85, 105)
                        else:
                            run_d = p_b.add_run()
                            run_d.text = b_text_clean
                            run_d.font.bold = False
                            run_d.font.color.rgb = RGBColor(71, 85, 105)
                            
                else:
                    # Layout C: CALLOUT_CARDS (Pila de Tarjetas Ejecutivas a todo el ancho)
                    bullets_callout = slide_data.get("bullets", [])
                    N = max(1, len(bullets_callout))
                    if N == 1:
                        card_height = Inches(1.8)
                        gap = Inches(0)
                    elif N == 2:
                        card_height = Inches(1.6)
                        gap = Inches(0.4)
                    elif N == 3:
                        card_height = Inches(1.2)
                        gap = Inches(0.3)
                    else:
                        card_height = Inches(0.9)
                        gap = Inches(0.15)
                    
                    for i, b_text in enumerate(bullets_callout):
                        card_top = Inches(2.0) + i * (card_height + gap)
                        
                        # 1. Fondo de la tarjeta (Slate 100)
                        bg_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), card_top, Inches(12.33), card_height)
                        bg_card.fill.solid()
                        bg_card.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate 100
                        bg_card.line.fill.background()
                        
                        # 2. Barra de acento a la izquierda (Azul Vektra)
                        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), card_top, Inches(0.08), card_height)
                        accent_bar.fill.solid()
                        accent_bar.fill.fore_color.rgb = RGBColor(37, 99, 235) # Azul Vektra
                        accent_bar.line.fill.background()
                        
                        # 3. Cuadro de texto superpuesto
                        txBox = slide.shapes.add_textbox(Inches(0.5 + 0.25), card_top + Inches(0.05), Inches(12.33 - 0.4), card_height - Inches(0.1))
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        tf.margin_left = Inches(0)
                        tf.margin_top = Inches(0)
                        tf.margin_right = Inches(0)
                        tf.margin_bottom = Inches(0)
                        
                        p = tf.paragraphs[0]
                        font_size = Pt(13) if N <= 3 else Pt(11.5)
                        p.font.name = "Arial"
                        p.font.size = font_size
                        
                        # Negrita inteligente
                        b_text_clean = clean_text(b_text).strip()
                        parts = b_text_clean.split(":", 1)
                        if len(parts) == 2 and len(parts[0]) < 40:
                            run_title = p.add_run()
                            run_title.text = parts[0].strip() + ": "
                            run_title.font.bold = True
                            run_title.font.size = font_size
                            run_title.font.color.rgb = RGBColor(30, 41, 59) # Slate 800
                            
                            run_desc = p.add_run()
                            run_desc.text = parts[1].strip()
                            run_desc.font.bold = False
                            run_desc.font.size = font_size
                            run_desc.font.color.rgb = RGBColor(71, 85, 105) # Slate 600
                        else:
                            run_all = p.add_run()
                            run_all.text = b_text_clean
                            run_all.font.bold = False
                            run_all.font.size = font_size
                            run_all.font.color.rgb = RGBColor(71, 85, 105) # Slate 600
            else:
                # Layout con gráfico: Texto estructurado a la izquierda, gráfico a la derecha
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(4.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                
                font_size = Pt(13) if len(bullets) <= 3 else Pt(11.5)
                
                first_p = True
                for b_text in bullets:
                    if first_p:
                        p = tf.paragraphs[0]
                        first_p = False
                    else:
                        p = tf.add_paragraph()
                        
                    p.font.name = "Arial"
                    p.font.size = font_size
                    p.space_after = Pt(8)
                    
                    # Añadir el punto de viñeta al inicio
                    run_bullet = p.add_run()
                    run_bullet.text = "• "
                    run_bullet.font.bold = False
                    run_bullet.font.size = font_size
                    run_bullet.font.color.rgb = RGBColor(75, 85, 99)
                    
                    # Negrita inteligente
                    b_text_clean = clean_text(b_text).strip()
                    parts = b_text_clean.split(":", 1)
                    if len(parts) == 2 and len(parts[0]) < 40:
                        run_title = p.add_run()
                        run_title.text = parts[0].strip() + ": "
                        run_title.font.bold = True
                        run_title.font.size = font_size
                        run_title.font.color.rgb = RGBColor(30, 41, 59) # Slate 800
                        
                        run_desc = p.add_run()
                        run_desc.text = parts[1].strip()
                        run_desc.font.bold = False
                        run_desc.font.size = font_size
                        run_desc.font.color.rgb = RGBColor(71, 85, 105) # Slate 600
                    else:
                        run_desc = p.add_run()
                        run_desc.text = b_text_clean
                        run_desc.font.bold = False
                        run_desc.font.size = font_size
                        run_desc.font.color.rgb = RGBColor(71, 85, 105) # Slate 600
                
            # Insertar Gráfico a la derecha (balanceado vertical y horizontalmente)
            if fig_data:
                img_bytes = export_plotly_to_image(fig_data)
                if img_bytes:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(img_bytes)
                        tmp_path = tmp.name
                    try:
                        # Insertar imagen alineada con el texto en posición left=5.3 (ancho 7.5)
                        slide.shapes.add_picture(tmp_path, Inches(5.3), Inches(2.0), width=Inches(7.5))
                    finally:
                        if os.path.exists(tmp_path):
                            os.remove(tmp_path)

        prs.save(output_path)
        return True, output_path
    except Exception as e:
        is_render = os.getenv("RENDER", "false").lower() == "true"
        if not is_render:
            logger.exception("Fallo en generación PPTX: %s", type(e).__name__)
        else:
            logger.error("Fallo en generación PPTX: %s", type(e).__name__)
        return False, f"Fallo en generación PPTX: {type(e).__name__}"

